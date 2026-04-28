"""Build the Phase 2 PowerPoint deck and Word report into deliverables/.

Re-runnable: `python3 build_deliverables.py`. Requires charts/ to already exist
(run phase2_charts.py first if needed).
"""
from __future__ import annotations

import os
from datetime import date

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches as DocxInches
from docx.shared import Pt as DocxPt
from docx.shared import RGBColor as DocxRGB
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ------------------------------------------------------------------
# Branding
# ------------------------------------------------------------------
NAVY = RGBColor(0x01, 0x21, 0x69)
GOLD = RGBColor(0xF2, 0xA9, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x22, 0x22, 0x22)

DOCX_NAVY = DocxRGB(0x01, 0x21, 0x69)
DOCX_GOLD = DocxRGB(0xF2, 0xA9, 0x00)

TITLE = "LGS Master's Programs: Tuition Pricing & Scholarship Analysis"
SUBTITLE = "Emory Laney Graduate School — Phase 2 (v2)"
REPORT_DATE = date(2026, 4, 23).strftime("%B %d, %Y")

# ------------------------------------------------------------------
# Content (insights keyed to each chart)
# ------------------------------------------------------------------
SEMESTER_SUBTEXT = "df_pricing population (gross tuition ≥ $20k). Sorted by Fall gross tuition."

# Each entry is either a (insight, chart_path) tuple — built with the standard
# top-strip layout — or a {headline, subtext, chart} dict, built with a left
# navy panel layout (used for the three per-semester slides).
CHART_SLIDES = [
    ("Discount rates vary six-fold across programs — MDP tops the list at 55% while Data Science and MBID give nothing",
     "charts/01_discount_rate_by_program.png"),
    {
        "headline": "Fall billing is tight across programs — Economics discounts most at the per-semester level",
        "subtext":  SEMESTER_SUBTEXT,
        "chart":    "charts/gross_net_fall.png",
    },
    {
        "headline": "Spring mirrors Fall with one exception — CS 4+1 shows the largest proportional discount",
        "subtext":  SEMESTER_SUBTEXT,
        "chart":    "charts/gross_net_spring.png",
    },
    {
        "headline": "Summer is where MDP's scholarship strategy concentrates — near-zero net on $22k gross",
        "subtext":  SEMESTER_SUBTEXT,
        "chart":    "charts/gross_net_summer.png",
    },
    {
        "headline": "High international share doesn't come with lower aid — CS is 86% international and 49% discounted",
        "subtext":  "Full population (n=487). Bubble size = enrollment; x-axis on log scale to spread the cluster of programs near 0%.",
        "chart":    "charts/03_intl_pct_vs_discount_rate.png",
    },
    ("Computer Science enrollment has roughly doubled since 2022 — driving nearly all of Emory's LGS master's volume growth",
     "charts/04_enrollment_trend_by_program.png"),
    ("Discount rates are broadly stable year-over-year — no program is systematically tightening aid",
     "charts/05_discount_rate_trend_by_program.png"),
    ("Emory's annualized tuition ($47–51k) undercuts Columbia/NYU by $20–40k and sits above Georgia Tech's in-state rates — mid-market positioning",
     "charts/06_peer_benchmark_emory_vs_peers.png"),
]

# v2: pricing metrics below (avg gross/net tuition, discount rate, peer
# annualization) are computed on df_pricing — student-program rows with
# total billed tuition of at least $20,000. This drops 11 of 493 rows
# (partial-term or single-course billings) that were pulling the averages
# down and compressing discount rates. Enrollment counts and
# international-mix figures still use the full population.
TOP5 = [
    ("MDP is Emory's costliest and most-discounted master's",
     "$108k average gross, $49k net, 55% discount rate across 135 full-load students over the 8-year panel ($8.4M in LGS scholarships)."),
    ("Two programs operate at 0% discount",
     "Data Science (n=11) and MBID (n=8) show gross = net for every full-load student. Obvious candidates for a pricing review."),
    ("Enrollment is bimodal",
     "Computer Science (154) and MDP (135) together account for 59% of the 493-student sample; eight other programs share the remaining 41%."),
    ("Computer Science is 86% international — and still ~49% discounted",
     "Emory competes for these students on price as well as visa-friendliness; aid does not scale down as international share rises."),
    ("Emory sits mid-pack against peers",
     "Annualized gross tuition ~$47–51k across Data/CS, Economics, and General buckets — below Columbia/NYU ($65–91k) and above Georgia Tech's in-state rates ($31–41k)."),
]

RECOMMENDATIONS = [
    ("Audit Data Science & MBID scholarship eligibility",
     "Both programs show 0% kept LGS scholarship across the 8-year panel "
     "(Data Science n=11, MBID n=8). Reconcile the raw FA ledger against the "
     "filtered scholarship feed to confirm whether this is a data anomaly or "
     "a deliberate full-pay carve-out — and if deliberate, decide whether it "
     "is costing yield against discounted peers."),
    ("Review CS aid strategy relative to international demand",
     "Computer Science is 86% international and still discounts ~49% per "
     "student — the highest international share and one of the deepest "
     "discounts in the portfolio. Test whether tighter aid materially shifts "
     "yield given that Columbia/NYU charge $20–40k more in absolute terms."),
    ("Capitalize on the peer pricing gap",
     "Annualized gross tuition ($47–51k) sits well below Columbia/NYU "
     "($65–91k) and clearly above Georgia Tech in-state ($31–41k). With CS "
     "enrollment roughly doubling since 2022, demand is responding — modeling "
     "selective sticker-price increases in the most-demanded programs (CS, "
     "MDP) is the lowest-risk revenue lever."),
    ("Stabilize MDP discount rate with a formal aid policy",
     "MDP discounts at 55% career-wide and concentrates the give in Summer "
     "(~$22k gross → ~$2k net per Summer term). Year-to-year the rate drifts "
     "in a 55–60% band, which suggests the practice is implicit. Codify the "
     "rule (eligibility, award size, term distribution) so future cohorts are "
     "predictable and the rate can be tuned deliberately rather than by drift."),
]

APPENDIX_BULLETS = [
    "Source data: institutional Excel exports of gross tuition billed and LGS scholarship awards; 8-year panel spanning AY 2019 through AY 2026.",
    "Net tuition is clipped at zero. In rare cases a kept scholarship exceeds billed tuition for a partial-term student; we report that student as paying $0 rather than a negative amount.",
    "Scholarship filter retains descriptions containing 'lgs' and drops external/federal awards (Pell, GRFP, Yellow Ribbon, Americorps, etc.) plus generic LGS 'special' awards, so discount rates reflect institutional aid only.",
    "Small-n programs are merged: ECONMS + ECON4P1MS → Economics; DATASCIMS + QTMMS → Data Science; DEVPRACMDP + HUMANRTCRT → MDP.",
    "Annualization for peer comparison uses program-specific terms-per-year (from the Information Sheet plus empirical seasonal distribution) rather than a blanket 2-term assumption. This corrects a ~33% understatement for 3-term programs such as MDP and Data Science.",
    "Programs with zero kept scholarship (Data Science, MBID) are retained in the summary with a 0% discount rate — this is intentional and flags potential pricing discrepancies.",
    "v2 pricing filter (df_pricing): avg gross/net tuition, discount rate, and peer benchmark exclude 11 of 493 student-program rows with total billed tuition under $20,000. These reflect single-course or partial-term billings with ~$0 scholarship that were depressing averages and compressing the gross-vs-net gap. Enrollment counts, international mix, and per-year trend charts continue to use the full population.",
]

REPORT_SECTIONS = [
    ("Discount Rates",
     ("Discount rates range from 0% to 55% across Emory's LGS master's programs. MDP, Computer Science, "
      "and Computer Science 4+1 all discount roughly half of sticker price, while Data Science and MBID "
      "give no institutional scholarship at all. The 8-year trend is broadly stable — year-over-year "
      "fluctuations sit within normal cohort noise and no program is systematically tightening aid."),
     ["charts/01_discount_rate_by_program.png", "charts/05_discount_rate_trend_by_program.png"]),
    {
        "title": "Gross vs Net Tuition",
        "intro": (
            "Splitting the tuition story by semester clarifies how each program prices and aids "
            "individual terms. Fall and Spring sit in the same range — full-load programs bill "
            "roughly $20–24k gross per term — but Summer is sparser, more concentrated, and "
            "reveals the heaviest discounting strategy in the portfolio. All three charts use the "
            "df_pricing population (career-total tuition ≥ $20k) and share the same Fall-gross-"
            "descending program order so cross-semester comparison is direct."
        ),
        "subsections": [
            ("2a — Fall",
             ("All ten programs enroll in Fall. Gross billings cluster around $20–24k, with MBID "
              "and Data Science at the top showing zero institutional aid (gross = net). Economics "
              "shows the largest per-semester Fall discount among high-gross programs — $24k gross "
              "→ $16k net. Computer Science, CS 4+1, Math, and MDP each take roughly half off "
              "through LGS aid."),
             "charts/gross_net_fall.png"),
            ("2b — Spring",
             ("Spring closely mirrors Fall in both program order and gross magnitudes. The notable "
              "exception is Computer Science 4+1, where Spring billings drop to ~$19k gross and net "
              "falls to ~$9k — a ~52% per-semester discount, the largest proportional Spring give "
              "in the portfolio. Computer Science shows a similar Spring drop. Bioethics 4+1 "
              "collects nearly all of its $22k Spring sticker."),
             "charts/gross_net_spring.png"),
            ("2c — Summer",
             ("Summer is structurally different — only eight of ten programs enroll students; MBID "
              "and Economics are Fall-Spring only. The story is MDP: $22k gross → $2.3k net per "
              "student, by far the heaviest discounting in the LGS portfolio and where the program "
              "concentrates its scholarship spend. Cancer Biology 4+1 shows a similar pattern at "
              "lower gross. Data Science Summer gross equals net, consistent with its 0% "
              "institutional-aid finding portfolio-wide."),
             "charts/gross_net_summer.png"),
        ],
    },
    ("International Mix",
     ("Bubble size shows program enrollment; position shows international share vs discount rate. Computer "
      "Science is the clear outlier — 86% international yet still 49% discounted. Across the portfolio there "
      "is no inverse relationship between international share and aid: high-international programs are not "
      "systematically less aided, which contradicts the common assumption that international students "
      "primarily subsidize domestic ones."),
     ["charts/03_intl_pct_vs_discount_rate.png"]),
    ("Enrollment Trends",
     ("The 8-year panel (AY 2019–2026) shows Computer Science climbing from the mid-30s to 60+ students per "
      "year, accounting for most of the overall growth in LGS master's enrollment. MDP remains Emory's other "
      "large program, steady in the 30–40 range per year. Smaller programs (Data Science, MBID, Economics) "
      "sit below 10 students per year throughout the window."),
     ["charts/04_enrollment_trend_by_program.png"]),
    ("Peer Benchmarking",
     ("Using the corrected program-specific annualization applied to full-load students, Emory's gross "
      "tuition sits firmly mid-pack: $47,353 Data/CS, $46,719 Economics, $50,654 General. Columbia leads "
      "the peer set at $65–91k depending on discipline, with NYU close behind at ~$70–76k. Georgia Tech's "
      "in-state rates ($31–41k) undercut Emory in every bucket. The market splits into a high-tier "
      "(Columbia/NYU), a mid-tier (Emory), and a value tier (Georgia Tech)."),
     ["charts/06_peer_benchmark_emory_vs_peers.png"]),
    {
        "title": "Recommendations",
        "intro": (
            "Four pricing-strategy questions emerge from the analysis. Each is scoped to be "
            "actionable by LGS leadership in a single review cycle: a data audit, an aid-strategy "
            "review, a sticker-price test, and a policy codification."
        ),
        "subsections": [
            ("6.1 — Audit Data Science & MBID scholarship eligibility",
             "Both programs show 0% kept LGS scholarship across the 8-year panel (Data Science "
             "n=11, MBID n=8). Either the awards are not being made or they are landing under "
             "descriptions excluded by the LGS-only filter. A manual reconciliation against the "
             "raw FA ledger confirms which it is. If the absence of aid is genuine, decide "
             "whether it reflects an intentional full-pay positioning or an oversight that is "
             "costing yield against more aggressively discounted peer programs.",
             None),
            ("6.2 — Review CS aid strategy relative to international demand",
             "Computer Science is 86% international and still discounts ~49% per student — the "
             "highest international share and one of the deepest discounts in the LGS portfolio. "
             "There is no inverse relationship between international share and aid in the data, "
             "which contradicts the common assumption that international students subsidize "
             "domestic ones. The strategic question is whether the 49% discount is necessary to "
             "land these students given that Columbia and NYU charge $20–40k more in absolute "
             "terms; a yield-elasticity test on the next admit cycle would quantify the lift per "
             "aid dollar.",
             None),
            ("6.3 — Capitalize on the peer pricing gap",
             "Annualized gross tuition sits comfortably below Columbia and NYU ($65–91k) and "
             "clearly above Georgia Tech's in-state floor ($31–41k). With CS enrollment roughly "
             "doubling since 2022, demand for Emory's mid-tier positioning is responding. A "
             "selective sticker-price increase of $5–8k in the most-demanded programs (CS, MDP) "
             "would still leave Emory comfortably below the urban-elite tier and would generate "
             "meaningful incremental revenue without changing market position.",
             None),
            ("6.4 — Stabilize MDP discount rate with a formal aid policy",
             "MDP discounts at 55% career-wide — the highest in the portfolio — and the per-"
             "semester view shows the give is concentrated in Summer ($22k gross → $2.3k net). "
             "Year-to-year the program rate hovers in a 55–60% band, which suggests the practice "
             "is implicit rather than codified. Document the decision rule (eligibility, award "
             "size, term distribution) so future cohorts are predictable and so the rate can be "
             "tuned deliberately rather than by drift.",
             None),
        ],
    },
]

# ------------------------------------------------------------------
# Speaker notes (one section per slide in the deck)
# ------------------------------------------------------------------
SPEAKER_NOTES = [
    ("Slide 1 — Title",
     [
         "Scope: Emory LGS master's programs, AY 2019 through AY 2026, 10 programs, 493 student-program enrollments.",
         "This is the v2 cut: pricing metrics now exclude 11 partial-term rows under $20,000 in total billed tuition. Enrollment counts and the international mix still reflect the full 493.",
     ]),
    ("Slide 2 — Executive Summary",
     [
         "Lead with MDP: costliest program on both gross and net basis, and the single biggest absolute tuition give.",
         "Flag Data Science (n=11) and MBID (n=8) as 0%-discount anomalies — obvious targets for a pricing review.",
         "Bimodal enrollment: CS + MDP = 59% of the sample; this is a portfolio of two big programs plus a long tail.",
         "High international share doesn't bring aid down — CS is 86% international AND ~49% discounted.",
         "Emory is mid-pack: $47–51k annualized, below Columbia/NYU and above Georgia Tech's in-state rates.",
     ]),
    ("Slide 3 — Discount rate by program",
     [
         "Top: MDP 55%, CS 4+1 49.4%, Computer Science 49.2%. These three are the institutional-aid heavyweights.",
         "Middle band: Math 41%, Economics 34%, Cancer Biology 4+1 26%.",
         "Bottom: Bioethics 9.5% and Bioethics 4+1 16%; Data Science and MBID at 0%.",
         "Uses the pricing filter (482 full-load students). Dropping the 11 sub-$20k rows barely moves any individual bar — the comparisons are robust to that choice — but it keeps the metric defensible.",
     ]),
    {
        "title": "Slide 4 — Fall: Per-Semester Gross vs Net by Program",
        "what_to_say": [
            "Frame this as a per-semester view, not annualized — students are billed roughly twice these amounts across Fall + Spring.",
            "All ten programs appear in Fall. Gross billings cluster tightly in the $20–24k range; the right tail (Bioethics) sits lower at ~$14k.",
            "Top of the chart: MBID and Data Science show gross = net — zero LGS aid in Fall, consistent with their portfolio-wide 0% discount.",
            "Economics is the standout in Fall: $24k gross → $16k net, the largest absolute per-semester give among high-gross programs.",
            "Computer Science, CS 4+1, Math, and MDP each show roughly half off through LGS aid in Fall.",
        ],
        "tip": "Open by reminding the audience that the y-axis is per-semester. The same student is billed roughly twice this amount over Fall + Spring — and for 3-term programs (MDP, Data Science) there is a Summer billing on top of that.",
    },
    {
        "title": "Slide 5 — Spring: Per-Semester Gross vs Net by Program",
        "what_to_say": [
            "Spring closely mirrors Fall — same program order, similar gross magnitudes — except for Computer Science 4+1.",
            "CS 4+1: Spring gross drops to ~$19k while net falls to ~$9k — a ~52% Spring discount, the largest proportional Spring give in the portfolio.",
            "Computer Science also shows a measurable Spring softening of gross billing relative to Fall.",
            "Bioethics 4+1 collects close to its $22k Spring sticker — the smallest Spring discount among full-load programs.",
            "Same df_pricing filter; same Fall-descending program order to keep the eye anchored.",
        ],
        "tip": "If you're pressed for time you can fold Spring into Fall — the Spring chart is essentially a copy of Fall with one CS-side detail. The slide is included for completeness and because Spring is when the largest proportional CS 4+1 discount shows up.",
    },
    {
        "title": "Slide 6 — Summer: Per-Semester Gross vs Net by Program",
        "what_to_say": [
            "Summer is structurally different: only eight programs are present. MBID and Economics are Fall-Spring only and don't appear here.",
            "The headline: MDP. $22k gross, $2.3k net per Summer term — the heaviest discounting in the LGS portfolio.",
            "Cancer Biology 4+1 shows a similar pattern at lower gross ($17.5k → $1.7k).",
            "CS 4+1 is the only other program with a meaningful Summer discount; CS Summer billing drops to ~$3k for the small group that takes Summer terms.",
            "Data Science Summer gross ≈ net, consistent with the program's 0% LGS-aid finding portfolio-wide.",
        ],
        "tip": "MDP is your headline number here — pause on it. $22k → $2k per Summer is dramatic and worth letting land. If asked: yes, MDP students typically take three terms per year, so this isn't a stray summer billing — it's where roughly a third of MDP's discount dollars are spent.",
    },
    ("Slide 7 — International % vs Discount Rate",
     [
         "Bubble size = enrollment; x = international %, y = discount rate.",
         "CS is the big navy bubble upper-right: 86% international, 49% discount, 154 students.",
         "No downward relationship between international share and aid — contradicts the intuition that international students subsidize domestic ones.",
         "Uses the full population (no pricing filter) — we don't want to undercount international students on a count-based chart.",
     ]),
    ("Slide 8 — Enrollment Trends",
     [
         "Full panel, AY 2019 through AY 2026, full population.",
         "Computer Science climbs from the mid-30s into the 60s — roughly doubled since 2022 — this is the main volume story.",
         "MDP steady around 30–40 per year.",
         "Smaller programs (Data Science, MBID, Economics) stay below ~10 per year throughout.",
     ]),
    ("Slide 9 — Discount Rate Trend",
     [
         "By-program, year-over-year. Shapes tell you stability; no program is systematically tightening or loosening aid.",
         "MDP hovers at 55–60%. CS holds around 49%. 0%-discount programs stay flat at zero.",
         "Year-to-year wobble sits within cohort-size noise; don't over-read single-year moves.",
     ]),
    ("Slide 10 — Peer Benchmark",
     [
         "Emory's three program-group averages on pricing-filtered data: Data/CS $47,353, Economics $46,719, General $50,654.",
         "Columbia leads Data/CS at $64,800 and Economics at $90,732.",
         "NYU sits in the $70,000–$75,750 range across buckets.",
         "Georgia Tech in-state is the floor: $31,210–$41,390.",
         "Positioning: Emory is priced like a mid-tier private — well above a public flagship, well below the urban elites. That's a deliberate pricing decision, not an accident.",
         "Peer numbers are published sticker rates; Emory numbers use the pricing filter so the comparison is like-for-like against a typical full-load student.",
     ]),
    {
        "title": "Slide 11 — Recommendations",
        "what_to_say_paragraphs": [
            "Frame this slide as the bridge from analysis to action. We've spent eight slides describing what the data says — these four items are what we'd actually do about it. They're scoped so LGS leadership could action all four in a single review cycle: one data audit, one aid-strategy review, one sticker-price test, and one policy codification.",
            "Recommendation 01 is the lowest-effort, highest-clarity item: Data Science and MBID both show literal-zero LGS scholarship across the entire 8-year panel. That is either a data-feed problem (awards are being made under descriptions our filter excludes) or it's a real pricing carve-out. We don't yet know which. A manual reconciliation against the raw FA ledger is a few hours of work and either way it changes the conversation.",
            "Recommendation 02 is the biggest dollar lever in the portfolio. Computer Science is 86% international and still gets ~49% off. The intuition that international students subsidize domestic ones simply doesn't show up in this data. Worth a yield-elasticity test on the next admit cycle to find out how much aid is actually buying us in conversion.",
            "Recommendation 03 is the upside case. Emory's $47–51k annualized sits $20–40k below Columbia and NYU. CS enrollment has roughly doubled since 2022 — the demand signal is there. A modest sticker increase in CS and MDP keeps us comfortably below the urban-elite tier and generates real revenue without repositioning the brand.",
            "Recommendation 04 is governance, not pricing. MDP's 55% discount rate is real and it's stable, but it isn't written down anywhere. That's how rates drift. Document the rule and the rate becomes a deliberate decision rather than an artifact.",
        ],
        "key_bullets": [
            "01 — Data audit. Cheapest item, biggest fact-check value.",
            "02 — CS aid review is the largest dollar conversation in the portfolio.",
            "03 — Sticker-price test is the lowest-risk revenue lever; CS/MDP are the obvious candidates.",
            "04 — MDP policy codification is governance work, not a pricing change.",
        ],
        "tip": "If you have time for only one rec in the room, lead with 02 — the CS aid question is where the conversation will gravitate anyway, and it's where the largest dollars are. 01 is a quick credibility win to set up before the bigger discussions.",
    },
    ("Slide 12 — Data & Methodology",
     [
         "Walk through the pipeline if asked: institutional Excel exports → decode term codes → filter to LGS programs → keep institutional-only scholarships → collapse to student-program level.",
         "The sub-$20k pricing filter: 11 rows under $20,000 in total billed tuition were excluded from pricing metrics only. These are partial-term billings with ~$0 kept scholarship that distort averages.",
         "Enrollment counts, international mix, and the per-year trend charts continue to run on the full 493 rows.",
         "Scholarship filter strips federal/external awards (Pell, GRFP, Yellow Ribbon, etc.) so the discount rate reflects institutional aid only.",
     ]),
]


# ------------------------------------------------------------------
# PowerPoint helpers
# ------------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill_color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    return shape


def add_text(slide, x, y, w, h, text, *, color=DARK_TEXT, size=18,
             bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def build_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    sw, sh = prs.slide_width, prs.slide_height
    # Full-bleed navy background
    add_rect(slide, 0, 0, sw, sh, NAVY)
    # Gold accent bar
    add_rect(slide, Inches(1), Inches(3.0), Inches(1.5), Inches(0.1), GOLD)
    # Title
    add_text(slide, Inches(1), Inches(3.2), sw - Inches(2), Inches(2),
             TITLE, color=WHITE, size=40, bold=True, align=PP_ALIGN.LEFT)
    # Subtitle
    add_text(slide, Inches(1), Inches(5.0), sw - Inches(2), Inches(0.6),
             SUBTITLE, color=GOLD, size=20, bold=True, align=PP_ALIGN.LEFT)
    # Date
    add_text(slide, Inches(1), Inches(6.8), sw - Inches(2), Inches(0.4),
             REPORT_DATE, color=WHITE, size=14, align=PP_ALIGN.LEFT)


def build_exec_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width
    # Navy heading bar
    add_rect(slide, 0, 0, sw, Inches(1.0), NAVY)
    add_text(slide, Inches(0.5), Inches(0.2), sw - Inches(1), Inches(0.7),
             "Executive Summary — Top 5 Findings",
             color=WHITE, size=28, bold=True)
    # Gold accent under heading
    add_rect(slide, 0, Inches(1.0), sw, Inches(0.08), GOLD)

    y = Inches(1.4)
    line_h = Inches(1.1)
    for i, (headline, detail) in enumerate(TOP5, start=1):
        # Number
        add_text(slide, Inches(0.5), y, Inches(0.7), line_h,
                 f"{i}.", color=GOLD, size=28, bold=True, align=PP_ALIGN.CENTER)
        # Headline + detail
        box = slide.shapes.add_textbox(Inches(1.2), y, sw - Inches(1.7), line_h)
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = headline
        r1.font.name = "Calibri"
        r1.font.size = Pt(16)
        r1.font.bold = True
        r1.font.color.rgb = NAVY
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = detail
        r2.font.name = "Calibri"
        r2.font.size = Pt(13)
        r2.font.color.rgb = DARK_TEXT
        y += line_h


def build_chart_slide(prs, insight, chart_path, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    # Navy headline strip
    add_rect(slide, 0, 0, sw, Inches(1.3), NAVY)
    add_text(slide, Inches(0.5), Inches(0.25), sw - Inches(1), Inches(0.85),
             insight, color=WHITE, size=20, bold=True)
    # Gold accent under headline
    add_rect(slide, 0, Inches(1.3), sw, Inches(0.06), GOLD)
    # Chart centered below
    img_y = Inches(1.6)
    img_h = sh - img_y - Inches(0.4)
    slide.shapes.add_picture(
        chart_path,
        left=Inches(0.5),
        top=img_y,
        width=sw - Inches(1.0),
        height=img_h,
    )
    # Page number bottom right
    add_text(slide, sw - Inches(0.8), sh - Inches(0.35), Inches(0.6), Inches(0.3),
             str(page_num), color=NAVY, size=11, align=PP_ALIGN.RIGHT)


def build_left_panel_chart_slide(prs, headline, subtext, chart_path, page_num):
    """Slide with a full-height left navy panel (headline + subtext) and a
    single chart on the right, sized to fit while preserving its aspect."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height

    panel_w = Inches(4.5)
    add_rect(slide, 0, 0, panel_w, sh, NAVY)
    add_rect(slide, panel_w, 0, Inches(0.06), sh, GOLD)

    add_text(
        slide, Inches(0.4), Inches(1.0), panel_w - Inches(0.7), Inches(4.0),
        headline, color=WHITE, size=20, bold=True, align=PP_ALIGN.LEFT,
    )
    add_text(
        slide, Inches(0.4), Inches(5.4), panel_w - Inches(0.7), Inches(1.6),
        subtext, color=GOLD, size=13, align=PP_ALIGN.LEFT,
    )

    # Right-side chart, fit to the available area while preserving aspect.
    right_x = panel_w + Inches(0.25)
    right_w = sw - right_x - Inches(0.3)
    avail_top = Inches(0.3)
    avail_bot = Inches(0.5)
    avail_h = sh - avail_top - avail_bot
    with Image.open(chart_path) as im:
        img_w_px, img_h_px = im.size
    chart_h_natural = right_w * img_h_px / img_w_px
    if chart_h_natural <= avail_h:
        chart_w = right_w
        chart_h = chart_h_natural
    else:
        chart_h = avail_h
        chart_w = chart_h * img_w_px / img_h_px
    chart_x = right_x + (right_w - chart_w) / 2
    chart_y = avail_top + (avail_h - chart_h) / 2
    slide.shapes.add_picture(chart_path, left=chart_x, top=chart_y, width=chart_w, height=chart_h)

    add_text(
        slide, sw - Inches(0.8), sh - Inches(0.35), Inches(0.6), Inches(0.3),
        str(page_num), color=NAVY, size=11, align=PP_ALIGN.RIGHT,
    )


def build_recommendations_slide(prs, page_num):
    """Full-bleed navy slide with four gold-numbered recommendation badges."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height

    add_rect(slide, 0, 0, sw, sh, NAVY)
    add_text(slide, Inches(0.6), Inches(0.4), sw - Inches(1.2), Inches(0.7),
             "Recommendations", color=GOLD, size=30, bold=True)
    add_rect(slide, Inches(0.6), Inches(1.15), Inches(1.5), Inches(0.06), GOLD)

    block_h = Inches(1.4)
    start_y = Inches(1.55)
    badge_d = Inches(0.95)
    text_left = Inches(2.0)
    text_right_margin = Inches(0.6)

    for i, (headline, body) in enumerate(RECOMMENDATIONS):
        y = start_y + block_h * i

        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.15), badge_d, badge_d,
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = GOLD
        badge.line.fill.background()
        tf = badge.text_frame
        tf.margin_left = tf.margin_right = Inches(0.0)
        tf.margin_top = tf.margin_bottom = Inches(0.0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"{i + 1:02d}"
        run.font.name = "Calibri"
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = NAVY

        add_text(
            slide, text_left, y + Inches(0.1), sw - text_left - text_right_margin, Inches(0.45),
            headline, color=GOLD, size=16, bold=True,
        )
        add_text(
            slide, text_left, y + Inches(0.55), sw - text_left - text_right_margin, Inches(0.85),
            body, color=WHITE, size=11,
        )

    add_text(
        slide, sw - Inches(0.8), sh - Inches(0.35), Inches(0.6), Inches(0.3),
        str(page_num), color=GOLD, size=11, align=PP_ALIGN.RIGHT,
    )


def build_appendix_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width
    add_rect(slide, 0, 0, sw, Inches(1.0), NAVY)
    add_text(slide, Inches(0.5), Inches(0.2), sw - Inches(1), Inches(0.7),
             "Data Notes & Methodology", color=WHITE, size=26, bold=True)
    add_rect(slide, 0, Inches(1.0), sw, Inches(0.08), GOLD)

    y = Inches(1.4)
    for bullet in APPENDIX_BULLETS:
        box = slide.shapes.add_textbox(Inches(0.5), y, sw - Inches(1), Inches(0.9))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = "• "
        r1.font.name = "Calibri"
        r1.font.size = Pt(14)
        r1.font.bold = True
        r1.font.color.rgb = GOLD
        r2 = p.add_run()
        r2.text = bullet
        r2.font.name = "Calibri"
        r2.font.size = Pt(13)
        r2.font.color.rgb = DARK_TEXT
        y += Inches(0.85)


def build_pptx(path: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_title_slide(prs)
    build_exec_summary_slide(prs)
    for i, entry in enumerate(CHART_SLIDES, start=3):
        if isinstance(entry, dict):
            build_left_panel_chart_slide(
                prs, entry["headline"], entry["subtext"], entry["chart"], i,
            )
        else:
            insight, chart = entry
            build_chart_slide(prs, insight, chart, i)
    build_recommendations_slide(prs, page_num=3 + len(CHART_SLIDES))
    build_appendix_slide(prs)

    prs.save(path)
    print(f"  wrote {path}")


# ------------------------------------------------------------------
# Word helpers
# ------------------------------------------------------------------
def _shading(color_hex: str):
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), color_hex)
    return shd


def _page_break(doc):
    p = doc.add_paragraph()
    run = p.add_run()
    run.add_break(WD_BREAK.PAGE)


def add_gold_rule(doc):
    """Thin gold horizontal rule via a 1-row, 1-col table with bottom border."""
    table = doc.add_table(rows=1, cols=1)
    cell = table.cell(0, 0)
    cell.text = ""
    # Set bottom border to gold, remove others
    tcPr = cell._tc.get_or_add_tcPr()
    tcBorders = OxmlElement("w:tcBorders")
    for edge in ("top", "left", "right"):
        b = OxmlElement(f"w:{edge}")
        b.set(qn("w:val"), "nil")
        tcBorders.append(b)
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), "18")  # eighths of a point ≈ 2.25pt
    bottom.set(qn("w:color"), "F2A900")
    tcBorders.append(bottom)
    tcPr.append(tcBorders)
    # Zero row height so the rule is just the border
    for p in cell.paragraphs:
        p.paragraph_format.space_before = DocxPt(0)
        p.paragraph_format.space_after = DocxPt(0)


def add_heading(doc, text, level=1, size=20):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = DocxPt(12)
    p.paragraph_format.space_after = DocxPt(4)
    run = p.add_run(text)
    run.font.name = "Calibri"
    run.font.size = DocxPt(size)
    run.font.bold = True
    run.font.color.rgb = DOCX_NAVY


def add_body(doc, text, size=11):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = DocxPt(8)
    run = p.add_run(text)
    run.font.name = "Calibri"
    run.font.size = DocxPt(size)
    run.font.color.rgb = DocxRGB(0x22, 0x22, 0x22)


def add_numbered(doc, n, headline, detail):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = DocxPt(6)
    rn = p.add_run(f"{n}. ")
    rn.font.name = "Calibri"
    rn.font.size = DocxPt(11)
    rn.font.bold = True
    rn.font.color.rgb = DOCX_GOLD
    rh = p.add_run(headline + " — ")
    rh.font.name = "Calibri"
    rh.font.size = DocxPt(11)
    rh.font.bold = True
    rh.font.color.rgb = DOCX_NAVY
    rd = p.add_run(detail)
    rd.font.name = "Calibri"
    rd.font.size = DocxPt(11)
    rd.font.color.rgb = DocxRGB(0x22, 0x22, 0x22)


def add_subheading(doc, text, size=11):
    """Small inline label (e.g. 'What to say') above a bullet list."""
    p = doc.add_paragraph()
    p.paragraph_format.space_before = DocxPt(4)
    p.paragraph_format.space_after = DocxPt(2)
    run = p.add_run(text)
    run.font.name = "Calibri"
    run.font.size = DocxPt(size)
    run.font.bold = True
    run.font.color.rgb = DOCX_NAVY


def add_tip_box(doc, text):
    """Light-gold callout box for delivery tips."""
    table = doc.add_table(rows=1, cols=1)
    cell = table.cell(0, 0)
    tcPr = cell._tc.get_or_add_tcPr()
    tcPr.append(_shading("FFF4D6"))
    p = cell.paragraphs[0]
    p.paragraph_format.space_before = DocxPt(2)
    p.paragraph_format.space_after = DocxPt(2)
    r1 = p.add_run("TIP — ")
    r1.font.name = "Calibri"
    r1.font.size = DocxPt(11)
    r1.font.bold = True
    r1.font.color.rgb = DOCX_GOLD
    r2 = p.add_run(text)
    r2.font.name = "Calibri"
    r2.font.size = DocxPt(11)
    r2.font.color.rgb = DocxRGB(0x22, 0x22, 0x22)


def add_bullet(doc, text):
    p = doc.add_paragraph(style="List Bullet")
    p.paragraph_format.space_after = DocxPt(4)
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = DocxPt(11)
    run.font.color.rgb = DocxRGB(0x22, 0x22, 0x22)


def build_docx(path: str) -> None:
    doc = Document()

    # Global default font
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = DocxPt(11)

    # ---------- Cover page ----------
    for _ in range(6):
        doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(TITLE)
    r.font.name = "Calibri"
    r.font.size = DocxPt(28)
    r.font.bold = True
    r.font.color.rgb = DOCX_NAVY

    add_gold_rule(doc)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(SUBTITLE)
    r.font.name = "Calibri"
    r.font.size = DocxPt(16)
    r.font.color.rgb = DOCX_GOLD

    for _ in range(2):
        doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(REPORT_DATE)
    r.font.name = "Calibri"
    r.font.size = DocxPt(14)
    r.font.color.rgb = DocxRGB(0x44, 0x44, 0x44)

    _page_break(doc)

    # ---------- Executive Summary ----------
    add_heading(doc, "Executive Summary", size=22)
    add_gold_rule(doc)
    add_body(
        doc,
        "This report summarizes pricing, scholarship discounting, and peer benchmarking for Emory's "
        "Laney Graduate School master's programs, using an 8-year panel of student-level tuition and "
        "LGS-scholarship records (AY 2019–2026).",
    )
    for i, (headline, detail) in enumerate(TOP5, start=1):
        add_numbered(doc, i, headline, detail)

    _page_break(doc)

    # ---------- Sections ----------
    for entry in REPORT_SECTIONS:
        if isinstance(entry, dict):
            add_heading(doc, entry["title"], size=18)
            add_gold_rule(doc)
            add_body(doc, entry["intro"])
            for sub_title, sub_narrative, sub_chart in entry["subsections"]:
                add_heading(doc, sub_title, size=14)
                add_body(doc, sub_narrative)
                if sub_chart:
                    doc.add_picture(sub_chart, width=DocxInches(6.5))
                    cap = doc.paragraphs[-1]
                    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
            _page_break(doc)
        else:
            title, narrative, charts = entry
            add_heading(doc, title, size=18)
            add_gold_rule(doc)
            add_body(doc, narrative)
            for chart in charts:
                doc.add_picture(chart, width=DocxInches(6.5))
                cap = doc.paragraphs[-1]
                cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
            _page_break(doc)

    # ---------- Appendix ----------
    add_heading(doc, "Data & Methodology", size=18)
    add_gold_rule(doc)
    add_body(
        doc,
        "The analysis pipeline is implemented in analysis.ipynb and produces three Excel outputs "
        "(program_summary_output.xlsx, student_program_level_output.xlsx, and "
        "student_program_year_output.xlsx).",
    )
    for bullet in APPENDIX_BULLETS:
        add_bullet(doc, bullet)

    add_heading(doc, "Terms-Per-Year Lookup", size=14)
    rows = [
        ("COMPSCIMS, CS4P1MS, MATHMS, ECON4P1MS", "2", "Information Sheet (Fall-Spring)"),
        ("BIOETHMA, BIOETH4P1, BMIDMS", "2", "empirical (Fall-Spring dominant)"),
        ("ECONMS", "3", "Information Sheet (Summer-Fall-Spring)"),
        ("DATASCIMS", "3", "Information Sheet (Fall-Spring-Summer)"),
        ("QTMMS, DEVPRACMDP, HUMANRTCRT, BBS4P1MS", "3", "empirical (all three seasons ~equal)"),
    ]
    table = doc.add_table(rows=1 + len(rows), cols=3)
    table.style = "Light Grid Accent 1"
    hdr_cells = table.rows[0].cells
    for i, h in enumerate(("acad_plan", "terms/yr", "source")):
        hdr_cells[i].text = ""
        p = hdr_cells[i].paragraphs[0]
        r = p.add_run(h)
        r.font.bold = True
        r.font.color.rgb = DOCX_NAVY
    for ri, (a, b, c) in enumerate(rows, start=1):
        table.rows[ri].cells[0].text = a
        table.rows[ri].cells[1].text = b
        table.rows[ri].cells[2].text = c

    doc.save(path)
    print(f"  wrote {path}")


def build_speaker_notes_docx(path: str) -> None:
    """Per-slide talking points for the v2 deck."""
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = DocxPt(11)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Speaker Notes — LGS Master's Programs Pricing Analysis")
    r.font.name = "Calibri"
    r.font.size = DocxPt(22)
    r.font.bold = True
    r.font.color.rgb = DOCX_NAVY

    add_gold_rule(doc)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(f"{SUBTITLE} · {REPORT_DATE}")
    r.font.name = "Calibri"
    r.font.size = DocxPt(13)
    r.font.color.rgb = DOCX_GOLD

    doc.add_paragraph()
    add_body(
        doc,
        "One section per deck slide. v2 updates reflect the sub-$20k pricing "
        "filter applied to df_pricing: 11 partial-term rows excluded from "
        "discount-rate, gross-vs-net, and peer-benchmark metrics only. "
        "Full population retained for enrollment counts and international mix.",
    )

    for entry in SPEAKER_NOTES:
        if isinstance(entry, dict):
            add_heading(doc, entry["title"], size=15)
            add_subheading(doc, "What to say")
            if "what_to_say_paragraphs" in entry:
                for para in entry["what_to_say_paragraphs"]:
                    add_body(doc, para)
            else:
                for b in entry["what_to_say"]:
                    add_bullet(doc, b)
            if "key_bullets" in entry:
                add_subheading(doc, "Key points")
                for b in entry["key_bullets"]:
                    add_bullet(doc, b)
            add_tip_box(doc, entry["tip"])
        else:
            slide_title, bullets = entry
            add_heading(doc, slide_title, size=15)
            for b in bullets:
                add_bullet(doc, b)
        doc.add_paragraph()

    doc.save(path)
    print(f"  wrote {path}")


def main() -> None:
    os.makedirs("deliverables", exist_ok=True)
    build_pptx("deliverables/emory_grad_pricing_v2.pptx")
    build_docx("deliverables/emory_grad_pricing_report_v2.docx")
    build_speaker_notes_docx("deliverables/emory_grad_pricing_speaker_notes_v2.docx")


if __name__ == "__main__":
    main()
